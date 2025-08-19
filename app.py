import streamlit as st
import google.generativeai as genai
import json
import re
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import requests

class PPTGenerator:
    def __init__(self, api_key, model_name="gemini-2.5-pro"):
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel(model_name)

    def generate_content_outline_expert(self, topic, num_slides=5, debug=False):
        prompt = (
            f'Generate a professional PowerPoint presentation outline on the topic: "{topic}".'
            f'\n- Create exactly {num_slides} slides.'
            f'\n- Each slide should have a clear, descriptive title and 3-5 concise bullet points.'
            f'\n- Use formal, business-appropriate language suitable for a professional audience.'
            f'\n- Make sure bullet points are informative, specific, and free of jargon.'
            f'\n- Organize the outline logically with an introduction, main ideas, challenges, opportunities, and a conclusion.'
            f'\n- Return the result STRICTLY in JSON format as a list of objects with keys "title" and "content".'
            f'\n- Example format: [{{"title": "Slide 1 Title", "content": ["First bullet", "Second bullet"]}}, ...]'
        )
        try:
            response = self.model.generate_content(prompt)
            content = getattr(response, "text", None)
            if debug and content:
                st.subheader("Raw Gemini Output")
                st.code(content)
            if not content:
                return None
            match = re.search(r"\[.*\]", content, re.DOTALL)
            if match:
                try:
                    slides = json.loads(match.group(0))
                    return self._normalize_slides(slides)
                except json.JSONDecodeError:
                    st.error("Failed to parse JSON response.")
                    return None
            st.error("JSON format not found in model output.")
            return None
        except Exception as e:
            st.error(f"Error generating content: {e}")
            return None

    def generate_content_outline(self, topic, num_slides=5, debug=False):
        result = self.generate_content_outline_expert(topic, num_slides, debug)
        if result:
            st.success(f"Generated {len(result)} slides using Gemini")
            return result
        st.warning("Using fallback outline template...")
        return self._get_fallback_outline(topic, num_slides)

    def _get_fallback_outline(self, topic, num_slides):
        slides = [
            {"title": f"{topic} Overview", "content": [f"Intro to {topic}"]},
            {"title": "Key Concepts", "content": [f"Basic ideas of {topic}"]},
            {"title": "Challenges", "content": [f"Common issues in {topic}"]},
            {"title": "Opportunities", "content": [f"Future scope of {topic}"]},
            {"title": "Conclusion", "content": [f"Summary of {topic}"]},
        ]
        return slides[:num_slides]

    def _normalize_slides(self, slides):
        normalized = []
        for slide in slides:
            title = slide.get("title", "Untitled")
            content = slide.get("content", "")
            if isinstance(content, str):
                bullets = re.split(r"\n|•|-", content)
                bullets = [b.strip() for b in bullets if b.strip()]
            elif isinstance(content, list):
                bullets = [str(b).strip() for b in content if str(b).strip()]
            else:
                bullets = [str(content)]
            # 🔹 Clean + shorten bullets
            cleaned = []
            for b in bullets:
                b = re.sub(r"[*_`~]", "", b)  # remove formatting symbols
                if len(b) > 80:
                    b = b[:77].rsplit(" ", 1)[0] + "..."  # shorten long bullets
                cleaned.append(b)
            normalized.append({"title": title.strip(), "content": cleaned[:6]})  # max 6 bullets per slide
        return normalized

    def create_presentation(self, slides_data, use_images=False, pexels_api_key=None):
        prs = Presentation()
        blank_layout = prs.slide_layouts[5]
        for slide_info in slides_data:
            slide = prs.slides.add_slide(blank_layout)
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
            title_tf = title_box.text_frame
            p = title_tf.add_paragraph()
            p.text = slide_info["title"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            # Content box (narrower, avoids image overlap)
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.5), Inches(4))
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            for bullet in slide_info["content"]:
                p = content_tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
            # Optional image
            if use_images:
                img_stream = self._get_image(slide_info["title"], pexels_api_key)
                if img_stream:
                    slide.shapes.add_picture(img_stream, Inches(6.2), Inches(2.5), Inches(3), Inches(2))
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def _get_image(self, query, pexels_api_key):
        if not pexels_api_key:
            return None
        try:
            headers = {"Authorization": pexels_api_key}
            resp = requests.get(f"https://api.pexels.com/v1/search?query={query}&per_page=1", headers=headers)
            if resp.status_code == 200:
                data = resp.json()
                if data.get("photos"):
                    img_url = data["photos"][0]["src"]["medium"]
                    img_data = requests.get(img_url).content
                    return BytesIO(img_data)
        except Exception:
            return None
        return None


def main():
    # Page configuration with a modern look and central alignment and a nice theme color
    st.set_page_config(
        page_title="AI PowerPoint Generator",
        page_icon="🎨",
        layout="centered",
        initial_sidebar_state="expanded",
    )

    # Custom CSS styles for decoration and better layout
    st.markdown(
        """
        <style>
        /* Hide default Streamlit menu/footer */
        #MainMenu {visibility: hidden;}
        footer {visibility: hidden;}

        /* Background gradient */
        .reportview-container {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: #ffffff;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        }

        /* Title style */
        .stTitle {
            font-weight: 900 !important;
            font-size: 3.2rem !important;
            text-shadow: 0 0 8px rgba(0,0,0,0.4);
        }

        /* Subtitle style for intro text */
        .stText {
            font-size: 1.2rem !important;
            margin-bottom: 2rem;
        }

        /* Sidebar title styling */
        .sidebar .sidebar-content {
            background: #2c1e5c;
            color: #ddecef;
            padding: 1rem;
            border-radius: 10px;
            font-size: 1rem;
            font-weight: 600;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        }

        /* Button styling */
        div.stButton > button:first-child {
            background: #764ba2;
            color: white;
            font-size: 1.1rem;
            font-weight: 600;
            padding: 0.8rem 1.8rem;
            border-radius: 25px;
            border: none;
            box-shadow: 0 4px 15px rgba(118, 75, 162, 0.6);
            transition: background-color 0.3s ease;
            margin-top: 1rem;
            width: 100%;
        }
        div.stButton > button:first-child:hover {
            background: #5a357d;
            box-shadow: 0 6px 20px rgba(90, 53, 125, 0.9);
            cursor: pointer;
        }

        /* Input and select box enhancements */
        input, select {
            background-color: #4d3a83 !important;
            color: #eee !important;
            border-radius: 8px !important;
            border: 1px solid #7859b0 !important;
            padding: 0.4rem 0.6rem !important;
            font-weight: 600 !important;
        }
        input::placeholder {
            color: #b3a9d3 !important;
            font-style: italic;
        }

        /* Checkbox label */
        label[for^=use_images], label[for^=debug_mode] {
            color: #ddd;
            font-weight: 600;
        }

        /* Download button styling */
        div.stDownloadButton > button {
            background: linear-gradient(90deg, #89f7fe 0%, #66a6ff 100%);
            color: #061436 !important;
            font-weight: 700 !important;
            border-radius: 30px !important;
            padding: 0.8rem 2rem !important;
            box-shadow: 0 8px 20px rgba(102, 166, 255, 0.4);
            border: none !important;
            transition: background 0.3s ease;
            width: 100% !important;
        }
        div.stDownloadButton > button:hover {
            background: linear-gradient(90deg, #66a6ff 0%, #89f7fe 100%);
            cursor: pointer;
        }

        /* Code block styling for debug output to match theme */
        .stCodeBlock pre {
            background-color: #3a2f6d !important;
            color: #ccc !important;
            border-radius: 10px !important;
            font-size: 0.95rem !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    st.title("🎨 AI PowerPoint Generator")
    st.markdown("Generate professional presentations with AI quickly and beautifully.", unsafe_allow_html=True)

    with st.sidebar:
        st.header("⚙️ Configuration")
        gemini_api_key = st.text_input(
            "Gemini API Key",
            type="password",
            help="Required: Get from Google AI Studio",
        )
        pexels_api_key = st.text_input(
            "Pexels API Key (Optional)",
            type="password",
            help="Use for fetching related stock images",
        )
        model_options = [
            "gemini-2.5-pro",
            "gemini-2.5-flash",
            "gemini-1.5-pro",
        ]
        selected_model = st.selectbox("Select AI Model", options=model_options, index=0)

        st.markdown("---")
        st.markdown("### Usage Tips")
        st.info(
            "• Use Gemini 2.5 Pro for best quality\n"
            "• Gemini Flash for faster results\n"
            "• Include images to enhance slides visually"
        )

    # Main input section: centered with clean fields
    st.markdown("### 📋 Presentation Setup")
    topic = st.text_input("Enter your presentation topic", placeholder="E.g., The Future of Renewable Energy")
    num_slides = st.slider("Number of slides", 3, 15, 5)
    use_images = st.checkbox("Include images in slides", value=True)
    debug_mode = st.checkbox("Debug: Show raw Gemini output", value=False)

    generate_clicked = st.button("🚀 Generate Presentation")

    if generate_clicked:
        # Validation
        if not gemini_api_key:
            st.error("⚠️ Gemini API key is required to generate presentations.")
        elif not topic or topic.strip() == "":
            st.error("⚠️ Please enter a valid presentation topic.")
        else:
            try:
                with st.spinner("Generating presentation content..."):
                    ppt = PPTGenerator(gemini_api_key, selected_model)
                    slides = ppt.generate_content_outline(topic.strip(), num_slides, debug=debug_mode)
                    ppt_file = ppt.create_presentation(slides, use_images, pexels_api_key)

                st.success("Presentation generated successfully! 🎉")
                st.download_button(
                    label="💾 Download PowerPoint (.pptx)",
                    data=ppt_file,
                    file_name=f"{topic.strip().replace(' ', '_')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download-pptx",
                    help="Click to download your AI-generated presentation"
                )
            except Exception as e:
                st.error(f"⚠️ Application error: {e}")
                st.code(traceback.format_exc())
                st.info("Try refreshing the page or check your API keys.")

if __name__ == "__main__":
    main()